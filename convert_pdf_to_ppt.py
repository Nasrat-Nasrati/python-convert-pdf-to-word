import os
import pdfplumber
from pptx import Presentation
from pptx.util import Inches

# Get the PDF file path from user input
pdf_path = input("Enter the PDF file path (e.g., E:/Python Projects/pdf_converter/pdf data/week2 of html.pdf): ")
output_folder = "converted_files"

# Check if the PDF file exists
if not os.path.exists(pdf_path):
    print("Error: File not found")
    exit()

# Create output folder if it doesn't exist
os.makedirs(output_folder, exist_ok=True)

# Generate output PowerPoint file path
pptx_path = os.path.join(output_folder, os.path.basename(pdf_path).replace(".pdf", ".pptx"))

# Process the PDF and convert to PowerPoint
try:
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Open the PDF and extract text
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:  # Check if text exists
                # Add a slide with a blank layout
                slide_layout = prs.slide_layouts[6]  # Blank slide layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add a text box to the slide
                left = top = Inches(1)  # Position 1 inch from left and top
                width = Inches(8)       # 8 inches wide
                height = Inches(5)      # 5 inches tall
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                
                # Add the extracted text to the text box
                p = text_frame.add_paragraph()
                p.text = text
    
    # Save the PowerPoint file
    prs.save(pptx_path)
    print(f"PowerPoint file created successfully: {pptx_path}")
except Exception as e:
    print(f"Error: Unable to create PowerPoint file: {e}")